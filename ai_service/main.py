import os
import time
import shutil
import json
import re
import zipfile
from datetime import datetime
import xml.etree.ElementTree as ET
from typing import List, Dict, Optional
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

# Load .env FIRST so GROQ_API_KEY and OPENAI_API_KEY are available before any placeholder logic
try:
    from dotenv import load_dotenv
    load_dotenv(dotenv_path=os.path.join(os.path.dirname(__file__), '.env'), override=True)
    print("[AI Service] .env loaded successfully.")
except ImportError:
    print("[AI Service] WARNING: python-dotenv not installed. Install with: pip install python-dotenv")
import pdfplumber

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

try:
    from summary_evaluation.routes import eval_routes
    app.include_router(eval_routes.router)
except ImportError as e:
    print(f"Failed to load summary_evaluation router: {e}")



@app.get("/health")
async def health_check():
    return {
        "status": "ok",
        "service": "ai-service",
    }


@app.get("/ready")
async def readiness_check():
    openai_key = os.environ.get("OPENAI_API_KEY")
    groq_key = os.environ.get("GROQ_API_KEY")
    return {
        "status": "ready",
        "dependencies": {
            "vectors_available": VECTORS_AVAILABLE,
            "llm_available": LLM_AVAILABLE,
            "openai_key_configured": bool(openai_key) and openai_key != "sk-placeholder",
            "groq_key_configured": bool(groq_key) and groq_key != "gsk-placeholder",
            "groq_enabled": groq_llm_enabled(),
            "groq_model": os.environ.get("GROQ_MODEL", "llama-3.1-8b-instant")
        },
    }

if not os.environ.get("OPENAI_API_KEY"):
    # In production, load from .env or real environment variables
    os.environ["OPENAI_API_KEY"] = "sk-placeholder"
if not os.environ.get("GROQ_API_KEY"):
    os.environ["GROQ_API_KEY"] = "gsk-placeholder"
FAISS_INDEX_PATH = "faiss_index"

COURSE_FAISS_ROOT = "course_faiss_indexes"
COURSE_STYLE_ROOT = "course_profiles"
COURSE_CHUNKS_ROOT = "course_chunks"

os.makedirs(COURSE_FAISS_ROOT, exist_ok=True)
os.makedirs(COURSE_STYLE_ROOT, exist_ok=True)
os.makedirs(COURSE_CHUNKS_ROOT, exist_ok=True)

class ChatRequest(BaseModel):
    message: str
    professor_style: str = "Focus on practical examples, encourage students to read documentations, and explain concepts step-by-step."
    student_level: str = "intermediate"
    history: List[Dict[str, str]] = []

class ChatResponse(BaseModel):
    reply: str

class MultimodalChatResponse(BaseModel):
    reply: str
    transcript: str = ""
    extracted_text: str = ""
    attachments: List[Dict] = []

class PersonalizeRequest(BaseModel):
    quiz_scores: Dict[str, int]
    weak_topics: List[str]
    recent_interactions: List[str]

class PersonalizeResponse(BaseModel):
    next_best_topic: str
    topics_to_revise: List[str]
    practice_questions: List[str]
    adaptive_message: str

class EvalRequest(BaseModel):
    submission_text: str
    rubric: str

class CriterionResult(BaseModel):
    criterion: str
    weight: int
    score: int
    evidence: str
    rationale: str

class EvalResponse(BaseModel):
    score: str
    strengths: List[str]
    weaknesses: List[str]
    suggestions: List[str]
    explanation: str
    criterion_breakdown: List[CriterionResult] = []

class AssignmentEvalRequest(BaseModel):
    assignment_text: str
    submission_text: str
    rubric: str = ""

class AssignmentEvalResponse(BaseModel):
    total_score: int
    max_score: int = 100
    grade_label: str
    is_relevant: bool
    is_incomplete: bool
    score_breakdown: Dict[str, int]
    strengths: List[str]
    mistakes: List[str]
    missing_concepts: List[str]
    improvement_suggestions: List[str]
    summary: str
    detailed_feedback: str

# Robust Import Block
VECTORS_AVAILABLE = False
LLM_AVAILABLE = False

try:
    from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from langchain_openai import OpenAIEmbeddings
    from langchain_huggingface import HuggingFaceEmbeddings
    from langchain_community.vectorstores import FAISS
    from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
    from langchain.chains.combine_documents import create_stuff_documents_chain
    from langchain.chains import create_retrieval_chain
    from langchain_core.messages import HumanMessage, AIMessage
    from langchain_core.output_parsers import JsonOutputParser
    VECTORS_AVAILABLE = True
except ImportError as e:
    print(f"Warning: Vector / embedding integration unavailable ({e}). Fallback Mock AI active.")

try:
    from langchain_groq import ChatGroq
    LLM_AVAILABLE = True
except ImportError as e:
    print(f"Warning: Groq LLM integration unavailable ({e}). Mock AI will be used.")

IMAGE_TO_TEXT_AVAILABLE = False
try:
    from PIL import Image
    import pytesseract
    IMAGE_TO_TEXT_AVAILABLE = True
except ImportError:
    pass

def get_vector_store():
    if not VECTORS_AVAILABLE: return None
    embeddings = OpenAIEmbeddings()
    if os.path.exists(FAISS_INDEX_PATH):
        return FAISS.load_local(FAISS_INDEX_PATH, embeddings, allow_dangerous_deserialization=True)
    return None

@app.post("/upload")
async def upload_material(file: UploadFile = File(...)):
    if not file.filename.endswith(('.pdf', '.docx', '.pptx')):
        raise HTTPException(status_code=400, detail="Invalid file type. Send standard course materials.")

    file_location = f"temp_{file.filename}"
    with open(file_location, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    if not VECTORS_AVAILABLE:
        os.remove(file_location)
        return {"message": "[MOCK] Document uploaded safely. LangChain pipeline bypassed due to local dependencies."}

    documents = []
    try:
        if file.filename.endswith('.pdf'):
            loader = PyPDFLoader(file_location)
            documents = loader.load()
        elif file.filename.endswith('.docx'):
            loader = Docx2txtLoader(file_location)
            documents = loader.load()
            
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = text_splitter.split_documents(documents)
        
        embeddings = OpenAIEmbeddings()
        if os.path.exists(FAISS_INDEX_PATH):
            vector_store = FAISS.load_local(FAISS_INDEX_PATH, embeddings, allow_dangerous_deserialization=True)
            vector_store.add_documents(chunks)
        else:
            vector_store = FAISS.from_documents(chunks, embeddings)
            
        vector_store.save_local(FAISS_INDEX_PATH)
        
    except Exception as e:
        if os.path.exists(file_location): os.remove(file_location)
        raise HTTPException(status_code=500, detail=str(e))
        
    os.remove(file_location)
    return {"message": "Document processed, chunked, and stored into FAISS vector space successfully."}

@app.post("/chat", response_model=ChatResponse)
async def chat_with_agent(req: ChatRequest):
    if req.message == "dummy_test_connection":
        return ChatResponse(reply="Agent Microservice is ALIVE and connected.")
        
    if not VECTORS_AVAILABLE:
        # Graceful fallback logic so UI remains functional
        fallback_reply = "I don't have enough course material yet, please ask your professor to upload content."
        return ChatResponse(reply=fallback_reply)

    vector_store = get_vector_store()
    if not vector_store:
        return ChatResponse(reply="I am your Course Agent. Please upload the syllabus and course materials so I can study them for context!")

    retriever = vector_store.as_retriever(search_kwargs={"k": 4})
    
    # 2. Advanced Prompt Engineering
    system_prompt = (
        "You are a friendly teaching assistant. Answer the student's question concisely using the provided course material.\n\n"
        "Rules:\n"
        "1. Answer like a friendly teaching assistant.\n"
        "2. DO NOT include section titles like 'Course Material Response', 'Relevant Material', etc.\n"
        "3. Use a simple explanation based on course material.\n"
        "4. Keep response concise and natural (max 5-6 lines).\n"
        "5. If no relevant material is found in the context, say: 'I don’t have enough course material yet, please ask your professor to upload content.'\n"
        "6. Return ONLY the final answer string. Use plain text and avoid markdown headings.\n"
        f"7. Adaptation: Student level is '{req.student_level}'. Professor's style: '{req.professor_style}'.\n\n"
        "Context:\n{context}"
    )
    
    prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        MessagesPlaceholder(variable_name="chat_history"),
        ("human", "{input}")
    ])
    
    try:
        # Mock mode if Groq isn't available or API key is placeholder.
        if (not LLM_AVAILABLE) or os.environ.get("GROQ_API_KEY") == "gsk-placeholder":
            docs = retriever.invoke(req.message)
            source_content = docs[0].page_content if docs else "No specific documents found."
            
            # Clean conversational mock response
            mock_reply = f"Based on the course materials: {source_content[:500]}. I hope this helps you understand '{req.message}' better!"
            return ChatResponse(reply=mock_reply)

        # 3. Groq LLM (Llama 3 8b or Mixtral)
        llm = ChatGroq(model_name="llama3-8b-8192", temperature=0.7)
            
        # 4. Langchain Memory Generation
        chat_history_messages = []
        for msg in req.history:
            if msg.get("sender") == "user":
                chat_history_messages.append(HumanMessage(content=msg.get("text", "")))
            elif msg.get("sender") == "agent":
                chat_history_messages.append(AIMessage(content=msg.get("text", "")))

        question_answer_chain = create_stuff_documents_chain(llm, prompt)
        rag_chain = create_retrieval_chain(retriever, question_answer_chain)
        
        response = rag_chain.invoke({
            "input": req.message,
            "chat_history": chat_history_messages
        })
        return ChatResponse(reply=response["answer"])
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

############################################################
# Course-specific RAG (professor uploads -> per-course FAISS)
############################################################

class CourseChatRequest(BaseModel):
    course_key: str
    message: str
    professor_style: Optional[str] = None
    student_level: str = "intermediate"
    history: List[Dict[str, str]] = []

class WeeklyUpdateRequest(BaseModel):
    course_key: str
    week_label: str = "Weekly Update"
    new_topics: List[str] = []
    announcements: List[str] = []
    revised_expectations: List[str] = []
    update_text: str = ""

class QuizQuestion(BaseModel):
    questionText: str
    options: List[str]
    correctAnswer: int
    explanation: str = ""
    topic: str = ""

class QuizGenerateRequest(BaseModel):
    course_key: str
    count: int = 10
    difficulty: str = "medium"
    topics: List[str] = []
    instructions: str = ""

class QuizGenerateResponse(BaseModel):
    course_key: str
    questions: List[QuizQuestion]

def normalize_course_key(course_key: str) -> str:
    s = str(course_key).strip().lower()
    s = re.sub(r"[^a-z0-9\\-_]+", "-", s)
    s = re.sub(r"-+", "-", s).strip("-")
    return s or "default"

def openai_embeddings_enabled() -> bool:
    if not VECTORS_AVAILABLE:
        return False
    openai_key = os.environ.get("OPENAI_API_KEY")
    return bool(openai_key) and openai_key != "sk-placeholder"

def groq_llm_enabled() -> bool:
    if not LLM_AVAILABLE:
        return False
    groq_key = os.environ.get("GROQ_API_KEY")
    return bool(groq_key) and groq_key != "gsk-placeholder"

def course_style_path(course_key: str) -> str:
    return os.path.join(COURSE_STYLE_ROOT, f"{course_key}.json")

def course_chunks_path(course_key: str) -> str:
    return os.path.join(COURSE_CHUNKS_ROOT, f"{course_key}.json")

def course_faiss_path(course_key: str) -> str:
    return os.path.join(COURSE_FAISS_ROOT, course_key)

def load_course_style(course_key: str) -> str:
    try:
        p = course_style_path(course_key)
        if os.path.exists(p):
            with open(p, "r", encoding="utf-8") as f:
                data = json.load(f)
            return str(data.get("teaching_style", "") or "")
    except Exception:
        pass
    return ""

def save_course_style(course_key: str, teaching_style: str) -> None:
    p = course_style_path(course_key)
    with open(p, "w", encoding="utf-8") as f:
        json.dump({ "teaching_style": teaching_style or "" }, f)

def tokenize_for_scoring(text: str) -> List[str]:
    return re.findall(r"[a-zA-Z0-9]+", str(text).lower())

def split_text_into_chunks(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
    cleaned = re.sub(r"\\s+", " ", str(text)).strip()
    if not cleaned:
        return []

    # Use LangChain splitter when available; otherwise fallback to a simple sliding window.
    if VECTORS_AVAILABLE and "RecursiveCharacterTextSplitter" in globals():
        splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        return splitter.split_text(cleaned)

    step = max(1, chunk_size - chunk_overlap)
    chunks = []
    for i in range(0, len(cleaned), step):
        chunk = cleaned[i : i + chunk_size]
        if chunk:
            chunks.append(chunk)
    return chunks

def extract_text_from_file(file_location: str, filename: str) -> str:
    lower = filename.lower()
    if lower.endswith(".pdf"):
        texts = []
        with pdfplumber.open(file_location) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                if t:
                    texts.append(t)
        return "\\n".join(texts).strip()

    if lower.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(file_location)
            return "\\n".join([p.text for p in doc.paragraphs if p.text]).strip()
        except Exception:
            # Fallback: extract raw XML text runs (works without python-docx).
            try:
                with zipfile.ZipFile(file_location) as z:
                    xml_bytes = z.read("word/document.xml")
                root = ET.fromstring(xml_bytes)
                texts = []
                for node in root.iter():
                    # Most Word runs store text in w:t
                    if node.tag.endswith("}t") and node.text:
                        texts.append(node.text)
                cleaned = "\\n".join(texts).strip()
                if cleaned:
                    return cleaned
            except Exception:
                pass

            # Fallback to LangChain loader if available.
            if VECTORS_AVAILABLE and "Docx2txtLoader" in globals():
                loader = Docx2txtLoader(file_location)
                docs = loader.load()
                return "\\n".join([d.page_content for d in docs]).strip()
            return ""

    if lower.endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(file_location)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\\n".join(parts).strip()
        except Exception:
            # Fallback: extract slide XML texts directly from the pptx zip.
            try:
                with zipfile.ZipFile(file_location) as z:
                    slide_xml_files = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
                    parts = []
                    for name in slide_xml_files:
                        xml_bytes = z.read(name)
                        root = ET.fromstring(xml_bytes)
                        for node in root.iter():
                            # Text is commonly in a:t
                            if node.tag.endswith("}t") and node.text:
                                parts.append(node.text)
                    cleaned = "\\n".join(parts).strip()
                    if cleaned:
                        return cleaned
            except Exception:
                pass
            return ""

    if lower.endswith((".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp")):
        # Optional OCR support for image-based submissions/materials.
        try:
            import pytesseract
            from PIL import Image

            img = Image.open(file_location)
            return str(pytesseract.image_to_string(img) or "").strip()
        except Exception:
            return ""

    if lower.endswith(".zip"):
        # Best-effort extraction for supported file types inside zip bundles.
        extracted_parts = []
        try:
            with zipfile.ZipFile(file_location) as archive:
                for member in archive.namelist():
                    if member.endswith("/"):
                        continue
                    member_lower = member.lower()
                    if not member_lower.endswith((
                        ".txt", ".md", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".json", ".csv",
                        ".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"
                    )):
                        continue

                    temp_member_path = f"{file_location}__{re.sub(r'[^a-zA-Z0-9._-]+', '_', os.path.basename(member) or 'member')}"
                    with archive.open(member) as src, open(temp_member_path, "wb") as dst:
                        shutil.copyfileobj(src, dst)

                    try:
                        content = extract_text_from_file(temp_member_path, os.path.basename(member) or member)
                        if not content and member_lower.endswith((".txt", ".md", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".json", ".csv")):
                            with open(temp_member_path, "r", encoding="utf-8", errors="ignore") as f:
                                content = f.read()

                        if content:
                            extracted_parts.append(f"### {member}\n{content[:20000]}")
                    finally:
                        try:
                            os.remove(temp_member_path)
                        except Exception:
                            pass
        except Exception:
            return ""

        return "\n\n".join(extracted_parts).strip()

    if lower.endswith((
        ".txt", ".md", ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".c", ".cpp", ".json", ".csv", ".html", ".css"
    )):
        try:
            with open(file_location, "r", encoding="utf-8", errors="ignore") as f:
                return f.read().strip()
        except Exception:
            return ""

    if lower.endswith((".mp3", ".wav", ".m4a", ".webm", ".ogg", ".mp4")):
        # Audio transcription is optional; supported when openai-whisper is installed.
        try:
            import whisper
            model_name = os.environ.get("WHISPER_MODEL", "base")
            model = whisper.load_model(model_name)
            result = model.transcribe(file_location)
            return str(result.get("text", "")).strip()
        except Exception:
            return ""

    return ""

def append_chunks_offline(course_key: str, new_items: List[Dict[str, str]]) -> None:
    p = course_chunks_path(course_key)
    existing = []
    if os.path.exists(p):
        try:
            with open(p, "r", encoding="utf-8") as f:
                existing = json.load(f) or []
        except Exception:
            existing = []
    existing.extend(new_items)
    with open(p, "w", encoding="utf-8") as f:
        json.dump(existing, f)

def course_faiss_exists(course_key: str) -> bool:
    p = course_faiss_path(course_key)
    return os.path.exists(p) and len(os.listdir(p)) > 0

def add_chunks_to_course_faiss(course_key: str, chunks: List[str], metadatas: List[Dict]) -> None:
    if not openai_embeddings_enabled():
        return

    embeddings = OpenAIEmbeddings()
    path = course_faiss_path(course_key)

    if course_faiss_exists(course_key):
        vector_store = FAISS.load_local(path, embeddings, allow_dangerous_deserialization=True)
        vector_store.add_texts(chunks, metadatas=metadatas)
    else:
        vector_store = FAISS.from_texts(chunks, embeddings, metadatas=metadatas)

    vector_store.save_local(path)

def retrieve_course_context(course_key: str, query: str, k: int = 4) -> List[str]:
    course_key = normalize_course_key(course_key)

    # Prefer FAISS retrieval when embeddings are enabled.
    if openai_embeddings_enabled() and course_faiss_exists(course_key):
        try:
            embeddings = OpenAIEmbeddings()
            vector_store = FAISS.load_local(course_faiss_path(course_key), embeddings, allow_dangerous_deserialization=True)
            docs = vector_store.similarity_search(query, k=k)
            return [d.page_content for d in docs if getattr(d, "page_content", None)]
        except Exception:
            pass

    # Offline fallback: naive token overlap scoring.
    try:
        p = course_chunks_path(course_key)
        if not os.path.exists(p):
            return []
        with open(p, "r", encoding="utf-8") as f:
            items = json.load(f) or []

        q_tokens = set(tokenize_for_scoring(query))
        if not q_tokens:
            return [items[i].get("text", "") for i in range(min(k, len(items)))]

        scored = []
        for it in items:
            t = str(it.get("text", "") or "")
            chunk_tokens = set(tokenize_for_scoring(t))
            score = len(q_tokens.intersection(chunk_tokens))
            scored.append((score, t))
        scored.sort(key=lambda x: x[0], reverse=True)
        top = [t for s, t in scored[:k] if t]
        return top
    except Exception:
        return []

def parse_history_payload(history_value) -> List[Dict[str, str]]:
    if isinstance(history_value, list):
        source = history_value
    else:
        try:
            source = json.loads(history_value or "[]")
        except Exception:
            source = []

    history_messages = []
    for item in source:
        if not isinstance(item, dict):
            continue
        text = str(item.get("text") or item.get("content") or "").strip()
        if not text:
            continue
        history_messages.append({
            "sender": "agent" if item.get("sender") == "agent" else "user",
            "text": text,
        })
    return history_messages[-40:]

def extract_attachment_records(files: List[UploadFile], temp_dir: str) -> List[Dict[str, str]]:
    records = []
    for uploaded in files or []:
        filename = uploaded.filename or "file"
        safe_name = re.sub(r"[^a-zA-Z0-9._-]+", "_", filename)
        file_location = os.path.join(temp_dir, safe_name)

        with open(file_location, "wb") as buffer:
            shutil.copyfileobj(uploaded.file, buffer)

        text = extract_text_from_file(file_location, filename)
        lower = filename.lower()
        if lower.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff")):
            file_type = "image"
        elif lower.endswith((".mp3", ".wav", ".m4a", ".webm", ".ogg")):
            file_type = "audio"
        else:
            file_type = "document"

        records.append({
            "filename": filename,
            "file_type": file_type,
            "text": text,
            "path": file_location,
        })
    return records

def build_attachment_context(attachment_records: List[Dict[str, str]]) -> str:
    sections = []
    for record in attachment_records or []:
        text = str(record.get("text") or "").strip()
        if not text:
            text = "No extractable text found. Use the attachment metadata and visual context if applicable."
        sections.append(
            f"[{record.get('file_type', 'file').upper()}] {record.get('filename', 'file')}\n{text[:1800]}"
        )
    return "\n\n".join(sections).strip()

def generate_course_chat_reply(
    course_key: str,
    message: str,
    history: List[Dict[str, str]],
    student_level: str = "intermediate",
    professor_style: Optional[str] = None,
    attachment_context: str = "",
) -> ChatResponse:
    course_key = normalize_course_key(course_key)
    style_from_store = load_course_style(course_key)
    teaching_style = str(professor_style or "").strip() or style_from_store or ""
    query = str(message or "").strip()
    if attachment_context:
        query = f"{query}\n\nUploaded attachment context:\n{attachment_context}"

    context_chunks = retrieve_course_context(course_key, query or message, k=4)
    context_text = "\n\n---\n\n".join(context_chunks).strip()

    if not context_text and attachment_context:
        context_chunks = [attachment_context]
        context_text = attachment_context

    if not context_text:
        return ChatResponse(
            reply="I don't have enough course material yet, please ask your professor to upload content."
        )

    if not groq_llm_enabled():
        top_chunk = context_chunks[0] if context_chunks else ""
        if not top_chunk:
            return ChatResponse(reply="I don't have enough course material yet, please ask your professor to upload content.")
        
        # Clean conversational fallback response
        reply = f"According to our course materials: {top_chunk[:500]}. For {student_level} level students, this covers the essentials. Let me know if you need more details!"
        return ChatResponse(reply=reply)

    retriever = None
    if huggingface_embeddings_enabled() and course_faiss_exists(course_key):
        embeddings = HuggingFaceEmbeddings()
        vector_store = FAISS.load_local(course_faiss_path(course_key), embeddings, allow_dangerous_deserialization=True)
        retriever = vector_store.as_retriever(search_kwargs={"k": 4})

    system_prompt = (
        "You are a friendly teaching assistant. Answer the student's question concisely using the provided course material.\n\n"
        "Rules:\n"
        "1. Answer like a friendly teaching assistant.\n"
        "2. DO NOT include section titles like 'Course Material Response', 'Relevant Material', etc.\n"
        "3. Use a simple explanation based on course material.\n"
        "4. Keep response concise and natural (max 5-6 lines).\n"
        "5. If no relevant material is found in the context, say: 'I don’t have enough course material yet, please ask your professor to upload content.'\n"
        "6. Return ONLY the final answer string. Use plain text and avoid markdown headings.\n"
        f"7. Adaptation: Student level is '{student_level}'. Professor's style: '{teaching_style or 'natural'}'.\n\n"
        "Context:\n{context}"
    )

    prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        MessagesPlaceholder(variable_name="chat_history"),
        ("human", "{input}"),
    ])

    llm = ChatGroq(model_name="llama-3.1-8b-instant", temperature=0.7)

    chat_history_messages = []
    for msg in history or []:
        if VECTORS_AVAILABLE:
            if msg.get("sender") == "user":
                chat_history_messages.append(HumanMessage(content=msg.get("text", "")))
            elif msg.get("sender") == "agent":
                chat_history_messages.append(AIMessage(content=msg.get("text", "")))

    if not retriever:
        return ChatResponse(
            reply=(
                "I found the relevant course context, but I can't run FAISS retrieval in this environment. "
                "I'll still answer using the retrieved snippet.\n\n"
                f"Snippet: {context_chunks[0][:600]}..."
            )
        )

    rag_chain = (
        {
            "context": retriever,
            "input": RunnablePassthrough(),
            "chat_history": lambda _: chat_history_messages,
        }
        | prompt
        | llm
        | StrOutputParser()
    )

    response = rag_chain.invoke(message)
    return ChatResponse(reply=response)

def _extract_json_array(text: str) -> Optional[List[Dict]]:
    try:
        obj = json.loads(text)
        if isinstance(obj, list):
            return obj
        if isinstance(obj, dict) and isinstance(obj.get("questions"), list):
            return obj.get("questions")
    except Exception:
        pass

    # Best-effort: pull the first JSON array from the response
    try:
        start = text.find("[")
        end = text.rfind("]")
        if start != -1 and end != -1 and end > start:
            candidate = text[start : end + 1]
            obj = json.loads(candidate)
            if isinstance(obj, list):
                return obj
    except Exception:
        pass
    return None

def fallback_generate_quiz(course_key: str, context_chunks: List[str], count: int, topics: List[str]) -> QuizGenerateResponse:
    # Deterministic, editable quiz scaffold when LLM is unavailable.
    joined = "\n\n".join([c for c in (context_chunks or []) if c])
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", joined) if s.strip()]
    if not sentences:
        sentences = [joined[:250]] if joined else ["Course material context is not available."]

    def mutate(s: str, salt: int) -> str:
        # Produce a slightly altered distractor from the same sentence.
        words = s.split()
        if len(words) < 6:
            return s + " (paraphrase)"
        drop = max(1, min(3, (salt % 3) + 1))
        return " ".join(words[:-drop]) + " ..."

    questions: List[QuizQuestion] = []
    for i in range(max(1, min(25, int(count or 10)))):
        base = sentences[i % len(sentences)]
        base = base[:220]
        correct = base
        opts = [
            correct,
            mutate(base, i + 1),
            mutate(base, i + 2),
            mutate(base, i + 3),
        ]
        topic = topics[i % len(topics)] if topics else ""
        questions.append(QuizQuestion(
            questionText=f"Which statement is supported by the course material?",
            options=opts,
            correctAnswer=0,
            explanation="Correct option is taken directly from the course material excerpt (fallback mode).",
            topic=topic,
        ))

    return QuizGenerateResponse(course_key=course_key, questions=questions)

def fallback_generate_quiz_from_topics(course_key: str, count: int, topics: List[str], difficulty: str, instructions: str) -> QuizGenerateResponse:
    # Deterministic quiz generation from topics only (no course materials required).
    # This keeps the UI functional even when embeddings/materials are absent.
    safe_topics = [t for t in (topics or []) if str(t).strip()]
    if not safe_topics:
        safe_topics = ["core concepts"]

    count = max(1, min(25, int(count or 10)))
    difficulty = str(difficulty or "medium").strip().lower()
    if difficulty not in {"easy", "medium", "hard"}:
        difficulty = "medium"

    # Template bank (kept generic to avoid claiming course-specific facts)
    templates = [
        ("Which option best describes '{topic}'?", "definition"),
        ("Which is a common use-case for '{topic}'?", "use_case"),
        ("Which statement about '{topic}' is MOST accurate?", "accuracy"),
        ("Which choice is an example related to '{topic}'?", "example"),
    ]

    def opts_for(tag: str, topic: str) -> List[str]:
        topic_clean = str(topic).strip()
        generic_wrong = [
            f"An unrelated concept not tied to {topic_clean}",
            f"A definition that contradicts the idea of {topic_clean}",
            f"A statement that is too broad to specifically describe {topic_clean}",
        ]

        if tag == "definition":
            correct = f"A concept or technique commonly discussed under {topic_clean}"
        elif tag == "use_case":
            correct = f"Applying {topic_clean} to solve a practical problem in its domain"
        elif tag == "accuracy":
            correct = f"{topic_clean} has trade-offs and should be applied with context"
        else:
            correct = f"A scenario where {topic_clean} is directly relevant"

        return [correct, generic_wrong[0], generic_wrong[1], generic_wrong[2]]

    questions: List[QuizQuestion] = []
    for i in range(count):
        topic = safe_topics[i % len(safe_topics)]
        prompt, tag = templates[i % len(templates)]
        qtext = prompt.format(topic=topic)
        options = opts_for(tag, topic)
        explanation = (
            f"This question was generated from the topic '{topic}'. "
            f"{('Difficulty: ' + difficulty + '. ') if difficulty else ''}"
            f"{('Instruction considered: ' + instructions[:140]) if instructions else ''}"
        ).strip()
        questions.append(
            QuizQuestion(
                questionText=qtext,
                options=options,
                correctAnswer=0,
                explanation=explanation,
                topic=str(topic),
            )
        )

    return QuizGenerateResponse(course_key=course_key, questions=questions)

def generate_quiz_core(req: QuizGenerateRequest) -> QuizGenerateResponse:
    course_key = normalize_course_key(req.course_key)
    count = max(1, min(25, int(req.count or 10)))
    difficulty = str(req.difficulty or "medium").strip().lower()
    if difficulty not in {"easy", "medium", "hard"}:
        difficulty = "medium"
    topics = [str(t).strip() for t in (req.topics or []) if str(t).strip()]
    instructions = str(req.instructions or "").strip()

    query = " ".join([
        "Generate multiple-choice practice questions",
        f"difficulty {difficulty}",
        ("topics: " + ", ".join(topics)) if topics else "",
        instructions,
    ]).strip()

    context_chunks = retrieve_course_context(course_key, query, k=6)
    context_text = "\n\n---\n\n".join(context_chunks).strip()

    # If there's no course context, still generate from topics/instructions.
    if not context_text:
        if not groq_llm_enabled():
            return fallback_generate_quiz_from_topics(course_key, count, topics, difficulty, instructions)
    else:
        if not groq_llm_enabled():
            return fallback_generate_quiz(course_key, context_chunks, count, topics)

    llm = ChatGroq(model_name="llama-3.1-8b-instant", temperature=0.35)
    parser = JsonOutputParser(pydantic_object=QuizGenerateResponse)

    prompt = ChatPromptTemplate.from_messages([
        ("system", (
            "You generate multiple-choice quizzes.\n"
            "If course excerpts are provided, prefer them; otherwise generate from the provided topics and general knowledge.\n"
            "Return ONLY valid JSON matching the schema below. Do not wrap in markdown.\n\n"
            "Course key: {course_key}\n"
            "Difficulty: {difficulty}\n"
            "Requested count: {count}\n"
            "Topics (optional): {topics}\n"
            "Extra instructions (optional): {instructions}\n\n"
            "Course excerpts (may be empty):\n{context}\n\n"
            "Rules:\n"
            "- Output must be JSON for schema: {format_instructions}\n"
            "- Provide exactly {count} questions.\n"
            "- Each question must have 4 options.\n"
            "- correctAnswer is 0-based index and must be within options.\n"
            "- Include a short explanation.\n"
            "- Avoid repeating the same concept verbatim across questions."
            "\n\nSafety/quality rules when course excerpts are empty:\n"
            "- Use only high-confidence, widely accepted definitions and concepts.\n"
            "- Avoid niche vendor-specific facts or tricky edge-case claims.\n"
            "- Keep explanations short and non-controversial.\n"
            "- Prefer conceptual questions grounded directly in the provided topics."
        )),
        ("human", "Generate the quiz now.")
    ])

    chain = prompt | llm | parser
    result = chain.invoke({
        "course_key": course_key,
        "difficulty": difficulty,
        "count": count,
        "topics": topics,
        "instructions": instructions,
        "context": context_text,
        "format_instructions": parser.get_format_instructions(),
    })

    try:
        obj = QuizGenerateResponse(**result)
    except Exception:
        # If parser returned malformed data, try to salvage from raw text
        raw = json.dumps(result) if not isinstance(result, str) else result
        arr = _extract_json_array(raw) or []
        obj = QuizGenerateResponse(course_key=course_key, questions=[QuizQuestion(**q) for q in arr])

    # Normalize length
    obj.questions = (obj.questions or [])[:count]
    return obj

@app.post("/course/generate-quiz", response_model=QuizGenerateResponse)
async def course_generate_quiz(req: QuizGenerateRequest):
    return generate_quiz_core(req)

def parse_rubric_criteria(rubric_text: str) -> List[Dict[str, int]]:
    rubric = str(rubric_text or "").strip()
    if not rubric:
        return [
            {"criterion": "Clarity", "weight": 30},
            {"criterion": "Technical Accuracy", "weight": 40},
            {"criterion": "Completeness", "weight": 30},
        ]

    # Supports patterns like: "Clarity (20%)" or "1. Accuracy 50%".
    matches = re.findall(r"(?:\d+[.)]\s*)?([^\n,;]+?)\s*\(?\s*(\d{1,3})\s*%\s*\)?", rubric)
    criteria = []
    for name, weight_str in matches:
        criterion = re.sub(r"\s+", " ", name).strip(" -:\t")
        if not criterion:
            continue
        weight = max(1, min(100, int(weight_str)))
        criteria.append({"criterion": criterion, "weight": weight})

    if not criteria:
        lines = [re.sub(r"\s+", " ", ln).strip() for ln in re.split(r"\n+", rubric) if ln.strip()]
        if not lines:
            lines = ["Clarity", "Technical Accuracy", "Completeness"]
        equal = max(1, int(round(100 / len(lines))))
        criteria = [{"criterion": ln[:80], "weight": equal} for ln in lines[:6]]

    total = sum(c["weight"] for c in criteria)
    if total <= 0:
        return [
            {"criterion": "Clarity", "weight": 30},
            {"criterion": "Technical Accuracy", "weight": 40},
            {"criterion": "Completeness", "weight": 30},
        ]

    # Normalize to exactly 100.
    normalized = []
    acc = 0
    for i, c in enumerate(criteria):
        if i == len(criteria) - 1:
            w = max(1, 100 - acc)
        else:
            w = max(1, int(round(c["weight"] * 100.0 / total)))
            acc += w
        normalized.append({"criterion": c["criterion"], "weight": w})

    # Fix possible overflow after rounding.
    overflow = sum(x["weight"] for x in normalized) - 100
    if overflow > 0:
        normalized[-1]["weight"] = max(1, normalized[-1]["weight"] - overflow)
    return normalized

def summarize_submission_signal(submission_text: str) -> Dict[str, int]:
    text = str(submission_text or "")
    words = re.findall(r"\w+", text)
    lines = text.splitlines()
    code_blocks = len(re.findall(r"```", text)) // 2
    headings = len(re.findall(r"^#{1,6}\s+", text, flags=re.MULTILINE))
    references = len(re.findall(r"\b(reference|citation|source|docs?)\b", text, flags=re.IGNORECASE))
    examples = len(re.findall(r"\b(example|for instance|e\.g\.)\b", text, flags=re.IGNORECASE))
    return {
        "word_count": len(words),
        "line_count": len(lines),
        "code_blocks": code_blocks,
        "headings": headings,
        "references": references,
        "examples": examples,
    }

def deterministic_anchor_score(criteria: List[Dict[str, int]], submission_text: str) -> int:
    signals = summarize_submission_signal(submission_text)
    quality = 0
    quality += 25 if signals["word_count"] >= 250 else 12 if signals["word_count"] >= 120 else 5
    quality += 20 if signals["code_blocks"] >= 2 else 10 if signals["code_blocks"] == 1 else 6
    quality += 15 if signals["headings"] >= 2 else 8 if signals["headings"] == 1 else 5
    quality += 20 if signals["references"] >= 2 else 10 if signals["references"] == 1 else 4
    quality += 20 if signals["examples"] >= 2 else 10 if signals["examples"] == 1 else 5
    return max(35, min(95, quality))

def fallback_eval(criteria: List[Dict[str, int]], submission_text: str) -> EvalResponse:
    anchor = deterministic_anchor_score(criteria, submission_text)
    breakdown = []
    for c in criteria:
        breakdown.append(CriterionResult(
            criterion=c["criterion"],
            weight=c["weight"],
            score=anchor,
            evidence="The submission shows clear effort and structure relevant to this criterion.",
            rationale="Reviewing the layout and content density indicates a good attempt at addressing the requirements.",
        ))

    return EvalResponse(
        score=f"{anchor}/100",
        strengths=["Submission has usable structure for rubric-based review."],
        weaknesses=["Detailed criterion-level semantic analysis is unavailable in fallback mode."],
        suggestions=["Add more explicit evidence per rubric criterion."],
        explanation="I've reviewed your submission's structure and content density. You've made a solid start—check the criterion breakdown for more specific feedback!",
        criterion_breakdown=breakdown,
    )

def coerce_score_0_100(value: object, default: int = 70) -> int:
    s = str(value or "").strip()
    frac = re.match(r"(\d+(?:\.\d+)?)\s*/\s*(\d+(?:\.\d+)?)", s)
    if frac:
        num = float(frac.group(1))
        den = float(frac.group(2))
        if den > 0:
            return max(0, min(100, int(round((num / den) * 100))))
    pct = re.match(r"(\d+(?:\.\d+)?)\s*%", s)
    if pct:
        return max(0, min(100, int(round(float(pct.group(1))))))
    try:
        n = int(round(float(s)))
        return max(0, min(100, n))
    except Exception:
        return max(0, min(100, int(default)))

def evaluate_submission_core(submission_text: str, rubric: str) -> EvalResponse:
    criteria = parse_rubric_criteria(rubric)

    if not LLM_AVAILABLE or os.environ.get("GROQ_API_KEY") == "gsk-placeholder":
        return fallback_eval(criteria, submission_text)

    llm = ChatGroq(model_name="llama3-8b-8192", temperature=0.0)
    parser = JsonOutputParser(pydantic_object=EvalResponse)
    anchor = deterministic_anchor_score(criteria, submission_text)

    prompt = ChatPromptTemplate.from_messages([
        ("system", (
            "You are an expert professor evaluating a student submission against a rubric.\n"
            "Your evaluation must be explainable, evidence-based, and deterministic in style.\n"
            "Do not use random scoring. Use the rubric weights exactly and provide criterion-wise reasoning.\n\n"
            "Rubric text:\n{rubric}\n\n"
            "Parsed weighted criteria (must follow):\n{criteria_json}\n\n"
            "Deterministic anchor score (structure-only baseline): {anchor}/100\n"
            "Your final weighted score may differ from the anchor by at most +/-10 unless there is explicit strong evidence.\n\n"
            "Return ONLY valid JSON for schema:\n"
            "{format_instructions}\n\n"
            "Scoring rules:\n"
            "1) Include criterion_breakdown entries for every criterion with criterion, weight, score, evidence, rationale.\n"
            "2) score in each criterion is 0-100 integer and must reflect rubric evidence.\n"
            "3) Overall score must equal weighted average of criterion scores.\n"
            "4) explanation should be natural and conversational (answering like a friendly professor). Avoid section titles or structured debug text.\n"
            "5) Keep the explanation concise (max 5-6 lines), use plain text, and avoid markdown headings (**)."
        )),
        ("human", "Submission to evaluate:\n{submission}")
    ])

    chain = prompt | llm | parser
    result = chain.invoke({
        "rubric": rubric,
        "criteria_json": json.dumps(criteria),
        "anchor": anchor,
        "submission": submission_text,
        "format_instructions": parser.get_format_instructions(),
    })

    eval_obj = EvalResponse(**result)

    if not eval_obj.criterion_breakdown:
        return fallback_eval(criteria, submission_text)

    # Normalize criterion weights and recompute final score deterministically.
    weighted_total = 0.0
    total_weight = 0
    cleaned_breakdown = []
    for idx, c in enumerate(eval_obj.criterion_breakdown):
        weight = c.weight if c.weight and c.weight > 0 else (criteria[idx]["weight"] if idx < len(criteria) else 1)
        score = coerce_score_0_100(c.score, default=anchor)
        weighted_total += score * weight
        total_weight += weight
        cleaned_breakdown.append(CriterionResult(
            criterion=c.criterion,
            weight=weight,
            score=score,
            evidence=c.evidence,
            rationale=c.rationale,
        ))

    if total_weight <= 0:
        final_score = anchor
    else:
        final_score = int(round(weighted_total / total_weight))

    return EvalResponse(
        score=f"{final_score}/100",
        strengths=eval_obj.strengths or ["Clear attempt to address rubric criteria."],
        weaknesses=eval_obj.weaknesses or ["Some rubric areas need stronger supporting evidence."],
        suggestions=eval_obj.suggestions or ["Revise based on criterion-level feedback before resubmitting."],
        explanation=eval_obj.explanation or "Evaluation computed from rubric-weighted criterion analysis.",
        criterion_breakdown=cleaned_breakdown,
    )

def _extract_key_topics(text: str, top_n: int = 12) -> List[str]:
    """Extract the most meaningful multi-word and single-word terms from text."""
    # Remove common stop-words for cleaner topic extraction
    stop = {
        'a','an','the','and','or','but','in','on','at','to','for','of','with','by',
        'is','was','are','were','be','been','being','have','has','had','do','does','did',
        'will','would','could','should','may','might','shall','can','it','its','this',
        'that','these','those','i','you','we','they','he','she','my','your','our','their',
        'from','as','not','no','so','if','then','each','all','any','more','also','how',
        'what','which','when','where','who','there','here','about','into','up','out',
        'use','used','using','get','set','make','made','just','only','very','well','both',
    }
    words = re.findall(r'[a-zA-Z]{3,}', text.lower())
    freq: Dict[str, int] = {}
    for w in words:
        if w not in stop:
            freq[w] = freq.get(w, 0) + 1
    sorted_words = sorted(freq.items(), key=lambda x: x[1], reverse=True)
    return [w for w, _ in sorted_words[:top_n]]


def build_assignment_fallback(assignment_text: str, submission_text: str) -> AssignmentEvalResponse:
    """Smart rule-based fallback that generates personalized, contextual feedback."""
    assignment_tokens = set(tokenize_for_scoring(assignment_text))
    submission_tokens = set(tokenize_for_scoring(submission_text))
    overlap = assignment_tokens.intersection(submission_tokens)
    submission_words = len(tokenize_for_scoring(submission_text))

    relevance_ratio = len(overlap) / max(1, len(assignment_tokens)) if assignment_tokens else 0.0
    is_relevant = relevance_ratio >= 0.1
    is_incomplete = submission_words < 120

    # Content signals
    signals = summarize_submission_signal(submission_text)
    has_code = signals['code_blocks'] > 0
    has_headings = signals['headings'] > 0
    has_examples = signals['examples'] > 0
    has_references = signals['references'] > 0

    # Score computation
    correctness = max(20, min(95, int(round(35 + relevance_ratio * 60))))
    topic_understanding = max(20, min(95, int(round(30 + relevance_ratio * 65))))
    completeness = max(15, min(95, int(round(min(submission_words, 600) / 6))))
    technical_accuracy = max(20, min(95, int(round((correctness + topic_understanding) / 2))))

    total = int(round((correctness + topic_understanding + completeness + technical_accuracy) / 4))
    if not is_relevant:
        total = min(total, 45)
    if is_incomplete:
        total = min(total, 55)

    grade_label = "A" if total >= 90 else "B" if total >= 80 else "C" if total >= 70 else "D" if total >= 60 else "F"

    # --- Smart strengths (based on what IS present) ---
    strengths: List[str] = []
    if is_relevant:
        covered = list(overlap)[:5]
        if covered:
            covered_str = ', '.join(covered[:3])
            strengths.append(f"Your submission covers key concepts from the assignment including: {covered_str}.")
        else:
            strengths.append("Your submission is generally aligned with the assignment topic.")
    if submission_words >= 300:
        strengths.append(f"Good depth — your response spans {submission_words} words, showing effort and elaboration.")
    elif submission_words >= 150:
        strengths.append("Adequate length with sufficient content to evaluate conceptual understanding.")
    if has_code:
        strengths.append("Includes code examples which demonstrates practical implementation knowledge.")
    if has_headings:
        strengths.append("Well-structured response with clear section organization.")
    if has_examples:
        strengths.append("Good use of examples to support concepts.")
    if has_references:
        strengths.append("Includes references or citations showing research effort.")
    if not strengths:
        strengths.append("Submission has been received and evaluated for all criteria.")

    # --- Smart weaknesses (what's missing or weak) ---
    missing_assignment_topics = list(assignment_tokens - submission_tokens)
    meaningful_missing = [
        t for t in missing_assignment_topics
        if len(t) > 4 and t not in {'please','write','explain','describe','create','implement','provide','include'}
    ][:5]

    mistakes: List[str] = []
    if not is_relevant:
        mistakes.append("The submission appears off-topic — many key assignment concepts are not addressed.")
    if is_incomplete:
        mistakes.append(f"Submission is too brief ({submission_words} words). A thorough response should be at least 200+ words.")
    if not has_code and 'implement' in assignment_text.lower():
        mistakes.append("Implementation details are missing — the assignment likely requires code or pseudocode.")
    if not has_examples:
        mistakes.append("No concrete examples provided to validate understanding.")
    if meaningful_missing:
        topic_str = ', '.join(meaningful_missing[:4])
        mistakes.append(f"Several assignment topics appear uncovered in the submission: {topic_str}.")
    if not mistakes:
        mistakes.append("Some rubric areas could benefit from more detailed supporting evidence.")

    # --- Missing concepts ---
    missing_concepts: List[str] = []
    assignment_topics = _extract_key_topics(assignment_text, top_n=10)
    submission_topics = set(_extract_key_topics(submission_text, top_n=20))
    for topic in assignment_topics:
        if topic not in submission_topics and len(topic) > 4:
            missing_concepts.append(f"Topic '{topic}' from assignment requirements not addressed in submission")
    if not missing_concepts:
        missing_concepts.append("Map each assignment requirement explicitly to a section of your response")

    # --- Improvement suggestions ---
    suggestions: List[str] = []
    if is_incomplete:
        suggestions.append("Expand your response — aim for at least 300-400 words covering all required topics.")
    if not has_code and 'implement' in assignment_text.lower():
        suggestions.append("Add code blocks with actual implementation, not just conceptual descriptions.")
    if not has_examples:
        suggestions.append("Include at least 2-3 concrete examples to demonstrate understanding.")
    if not has_headings:
        suggestions.append("Use section headings to organize your response clearly.")
    if meaningful_missing:
        suggestions.append(f"Revisit and address these topics: {', '.join(meaningful_missing[:3])}.")
    suggestions.append("Review the assignment rubric and ensure each criterion has explicit coverage in your submission.")

    # --- Natural summary sentence ---
    if total >= 80:
        summary = f"Strong submission (Grade {grade_label}) — your work covers most required concepts with good depth, though a few areas could be strengthened."
    elif total >= 65:
        summary = f"Good effort (Grade {grade_label}) — your submission covers the core ideas but lacks depth in some key areas. Expand on the missing topics for a higher score."
    elif total >= 50:
        summary = f"Partial submission (Grade {grade_label}) — several required concepts are present but the response needs significantly more detail and coverage."
    else:
        summary = f"Incomplete submission (Grade {grade_label}) — the response needs substantial revision to address the assignment requirements. Please review the missing concepts listed below."

    detailed_feedback = (
        f"Your submission was analyzed using semantic overlap and content structure analysis. "
        f"Out of {len(assignment_tokens)} key assignment terms, your submission covered {len(overlap)} ({int(relevance_ratio*100)}% relevance). "
        f"The response is {submission_words} words long with {signals['code_blocks']} code block(s), "
        f"{signals['headings']} section heading(s), and {signals['examples']} example(s). "
        f"Focus on addressing the missing concepts and expanding technical depth for a better score."
    )

    print(f"[Fallback Eval] Score={total} | Relevance={int(relevance_ratio*100)}% | Words={submission_words} | Strengths={len(strengths)} | Weaknesses={len(mistakes)}")

    return AssignmentEvalResponse(
        total_score=total,
        max_score=100,
        grade_label=grade_label,
        is_relevant=is_relevant,
        is_incomplete=is_incomplete,
        score_breakdown={
            "correctness": correctness,
            "topic_understanding": topic_understanding,
            "completeness": completeness,
            "technical_accuracy": technical_accuracy,
        },
        strengths=strengths,
        mistakes=mistakes,
        missing_concepts=missing_concepts[:5],
        improvement_suggestions=suggestions[:5],
        summary=summary,
        detailed_feedback=detailed_feedback,
    )

def _safe_parse_llm_json(raw: str) -> Dict:
    """Robustly extract a JSON object from an LLM response, even if wrapped in markdown fences."""
    # Strip markdown code fences if present
    cleaned = re.sub(r'^```(?:json)?\s*', '', raw.strip(), flags=re.IGNORECASE)
    cleaned = re.sub(r'```\s*$', '', cleaned.strip())
    try:
        return json.loads(cleaned)
    except Exception:
        pass
    # Try to find the first { ... } block
    start = cleaned.find('{')
    end = cleaned.rfind('}')
    if start != -1 and end > start:
        try:
            return json.loads(cleaned[start:end+1])
        except Exception:
            pass
    return {}


def evaluate_assignment_core(assignment_text: str, submission_text: str, rubric: str = "") -> AssignmentEvalResponse:
    print(f"[AssignEval] Starting evaluation | assignment_len={len(assignment_text)} | submission_len={len(submission_text)} | groq_enabled={groq_llm_enabled()}")
    print(f"[AssignEval] Submission preview: {submission_text[:200]}...")

    if not groq_llm_enabled():
        print("[AssignEval] Groq not enabled — using smart fallback.")
        return build_assignment_fallback(assignment_text, submission_text)

    print(f"[AssignEval] Groq enabled. Model=llama3-8b-8192. Sending to LLM...")

    llm = ChatGroq(model_name="llama3-8b-8192", temperature=0.0)
    parser = JsonOutputParser(pydantic_object=AssignmentEvalResponse)

    prompt = ChatPromptTemplate.from_messages([
        (
            "system",
            (
                "You are a strict but fair professor evaluating assignment submissions semantically.\n"
                "Evaluate if the submission is relevant, complete, conceptually correct, and technically accurate.\n"
                "Use the assignment context and rubric. Avoid keyword-only judgments; reason about concepts.\n"
                "Return ONLY valid JSON (no markdown fences, no extra text) using this schema:\n{format_instructions}\n"
                "IMPORTANT: All list fields (strengths, mistakes, missing_concepts, improvement_suggestions) must have at least 2 items each.\n"
                "The summary must be a complete, natural sentence explaining the student's performance."
            ),
        ),
        (
            "human",
            (
                "Assignment question/material:\n{assignment_text}\n\n"
                "Rubric:\n{rubric}\n\n"
                "Student submission:\n{submission_text}\n\n"
                "Evaluation rules:\n"
                "1) total_score must be a 0-100 integer reflecting actual submission quality.\n"
                "2) score_breakdown must have: correctness, topic_understanding, completeness, technical_accuracy (all 0-100).\n"
                "3) strengths: list what the student did WELL (at least 2 items).\n"
                "4) mistakes: list specific errors or gaps (at least 2 items).\n"
                "5) missing_concepts: list topics from the assignment that are absent in the submission (at least 2 items).\n"
                "6) improvement_suggestions: concrete actionable steps to improve (at least 2 items).\n"
                "7) summary: one clear sentence describing the overall quality.\n"
                "8) detailed_feedback: 3-5 sentences of qualitative analysis.\n"
                "9) Flag is_relevant (true if submission addresses the assignment topic).\n"
                "10) Flag is_incomplete (true if submission is too brief or missing major sections)."
            ),
        ),
    ])

    chain = prompt | llm | parser
    invoke_payload = {
        "assignment_text": assignment_text[:4000],
        "submission_text": submission_text[:4000],
        "rubric": rubric or "Evaluate holistically for correctness, completeness, and technical accuracy.",
        "format_instructions": parser.get_format_instructions(),
    }

    # Retry logic: up to 2 retries on failure
    result = None
    last_error = None
    for attempt in range(3):
        try:
            print(f"[AssignEval] LLM invoke attempt {attempt + 1}/3...")
            result = chain.invoke(invoke_payload)
            print(f"[AssignEval] LLM raw result type={type(result).__name__} | keys={list(result.keys()) if isinstance(result, dict) else 'N/A'}")
            break
        except Exception as e:
            last_error = e
            print(f"[AssignEval] Attempt {attempt + 1} failed: {e}")
            if attempt < 2:
                time.sleep(1.5)

    if result is None:
        print(f"[AssignEval] All LLM attempts failed ({last_error}). Using smart fallback.")
        return build_assignment_fallback(assignment_text, submission_text)

    # If parser returned a string (malformed), try to salvage JSON
    if isinstance(result, str):
        print(f"[AssignEval] Parser returned raw string — attempting JSON extraction.")
        result = _safe_parse_llm_json(result)

    # Validate and coerce all fields
    total_score = coerce_score_0_100(result.get("total_score", 0), default=60)
    breakdown = result.get("score_breakdown", {}) or {}

    def safe_list(val, default_msg: str) -> List[str]:
        lst = val if isinstance(val, list) else []
        lst = [str(x).strip() for x in lst if str(x).strip()]
        return lst if lst else [default_msg]

    normalized = AssignmentEvalResponse(
        total_score=total_score,
        max_score=100,
        grade_label=str(result.get("grade_label") or ("A" if total_score >= 90 else "B" if total_score >= 80 else "C" if total_score >= 70 else "D" if total_score >= 60 else "F")),
        is_relevant=bool(result.get("is_relevant", True)),
        is_incomplete=bool(result.get("is_incomplete", False)),
        score_breakdown={
            "correctness": coerce_score_0_100(breakdown.get("correctness", total_score), default=total_score),
            "topic_understanding": coerce_score_0_100(breakdown.get("topic_understanding", total_score), default=total_score),
            "completeness": coerce_score_0_100(breakdown.get("completeness", total_score), default=total_score),
            "technical_accuracy": coerce_score_0_100(breakdown.get("technical_accuracy", total_score), default=total_score),
        },
        strengths=safe_list(result.get("strengths"), "Shows understanding of core concepts."),
        mistakes=safe_list(result.get("mistakes"), "Some areas require further elaboration."),
        missing_concepts=safe_list(result.get("missing_concepts"), "Review assignment requirements for missing topics."),
        improvement_suggestions=safe_list(result.get("improvement_suggestions"), "Expand on key concepts with examples."),
        summary=str(result.get("summary") or f"Submission scored {total_score}/100 — review detailed feedback below."),
        detailed_feedback=str(result.get("detailed_feedback") or "Evaluation completed. See strengths and improvement areas above."),
    )

    print(f"[AssignEval] SUCCESS | score={normalized.total_score} | strengths={len(normalized.strengths)} | mistakes={len(normalized.mistakes)} | missing={len(normalized.missing_concepts)}")
    return normalized

@app.post("/course/upload")
async def upload_course_materials(
    course_key: str = Form(...),
    teaching_style: str = Form(""),
    files: List[UploadFile] = File(...),
):
    course_key = normalize_course_key(course_key)

    if not files:
        raise HTTPException(status_code=400, detail="No files provided")

    # Store teaching style so the assistant can mimic the professor.
    save_course_style(course_key, teaching_style)

    supported_exts = {".pdf", ".docx", ".pptx", ".mp3", ".wav", ".m4a"}
    for f in files:
        ext = os.path.splitext((f.filename or "").lower())[1]
        if ext not in supported_exts:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {f.filename}")

    processed_files = 0
    total_chunks = 0

    # Use a temp directory per request to avoid collisions.
    temp_dir = f"temp_course_upload_{course_key}"
    os.makedirs(temp_dir, exist_ok=True)

    try:
        for uploaded in files:
            filename = uploaded.filename or "file"
            safe_name = re.sub(r"[^a-zA-Z0-9._-]+", "_", filename)
            file_location = os.path.join(temp_dir, safe_name)

            with open(file_location, "wb") as buffer:
                shutil.copyfileobj(uploaded.file, buffer)

            text = extract_text_from_file(file_location, filename)
            if not text or len(text) < 30:
                raise HTTPException(status_code=400, detail=f"Could not extract enough text from {filename}")

            chunks = split_text_into_chunks(text)
            if not chunks:
                raise HTTPException(status_code=400, detail=f"No chunks produced for {filename}")

            metadatas = [
                {"source": filename, "course_key": course_key}
                for _ in range(len(chunks))
            ]

            append_chunks_offline(course_key, [{"text": c, "metadata": m} for c, m in zip(chunks, metadatas)])
            add_chunks_to_course_faiss(course_key, chunks, metadatas)

            processed_files += 1
            total_chunks += len(chunks)

    except Exception as e:
        print(f"[MaterialIngestion] CRITICAL ERROR during ingestion: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")
    finally:
        # Best-effort cleanup.
        try:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception:
            pass

    return {
        "message": "Course materials processed and stored successfully.",
        "course_key": course_key,
        "processed_files": processed_files,
        "chunks_added": total_chunks,
        "faiss_enabled": openai_embeddings_enabled(),
    }

@app.get("/course/list")
async def list_courses():
    try:
        keys = []
        if os.path.exists(COURSE_STYLE_ROOT):
            for name in os.listdir(COURSE_STYLE_ROOT):
                if name.endswith(".json"):
                    keys.append(os.path.splitext(name)[0])
        keys = sorted(list(set(keys)))
        return {"courses": keys}
    except Exception:
        return {"courses": []}

@app.post("/course/weekly-update")
async def ingest_weekly_update(req: WeeklyUpdateRequest):
    course_key = normalize_course_key(req.course_key)

    new_topics = [str(x).strip() for x in (req.new_topics or []) if str(x).strip()]
    announcements = [str(x).strip() for x in (req.announcements or []) if str(x).strip()]
    revised_expectations = [str(x).strip() for x in (req.revised_expectations or []) if str(x).strip()]
    update_text = str(req.update_text or "").strip()

    if not new_topics and not announcements and not revised_expectations and not update_text:
        raise HTTPException(status_code=400, detail="Weekly update content is required")

    blocks = [f"{req.week_label or 'Weekly Update'} - {course_key}"]
    if new_topics:
        blocks.append("New Topics:\n" + "\n".join([f"- {item}" for item in new_topics]))
    if announcements:
        blocks.append("Announcements:\n" + "\n".join([f"- {item}" for item in announcements]))
    if revised_expectations:
        blocks.append("Revised Expectations:\n" + "\n".join([f"- {item}" for item in revised_expectations]))
    if update_text:
        blocks.append("Additional Notes:\n" + update_text)

    merged_text = "\n\n".join(blocks).strip()
    chunks = split_text_into_chunks(merged_text, chunk_size=800, chunk_overlap=100)
    if not chunks:
        raise HTTPException(status_code=400, detail="No chunks produced from weekly update")

    ts = datetime.utcnow().isoformat() + "Z"
    metadatas = [
        {
            "source": "weekly-update",
            "type": "weekly_update",
            "course_key": course_key,
            "week_label": str(req.week_label or "Weekly Update"),
            "timestamp": ts,
        }
        for _ in range(len(chunks))
    ]

    append_chunks_offline(course_key, [{"text": c, "metadata": m} for c, m in zip(chunks, metadatas)])
    add_chunks_to_course_faiss(course_key, chunks, metadatas)

    return {
        "message": "Weekly update ingested successfully.",
        "course_key": course_key,
        "chunks_added": len(chunks),
        "faiss_enabled": openai_embeddings_enabled(),
    }

@app.post("/course/chat", response_model=ChatResponse)
async def course_chat(req: CourseChatRequest):
    print(f"[AI Service] Received /course/chat request: course={req.course_key}, message='{req.message}'")
    course_key = normalize_course_key(req.course_key)

    style_from_store = load_course_style(course_key)
    teaching_style = str(req.professor_style or "").strip() or style_from_store or ""

    context_chunks = retrieve_course_context(course_key, req.message, k=4)
    context_text = "\\n\\n---\\n\\n".join(context_chunks).strip()

    if not context_text:
        return ChatResponse(
            reply="I don't have enough course material yet, please ask your professor to upload content."
        )

    # If we can't use the LLM, generate a structured response from the course context
    if not groq_llm_enabled():
        top_chunk = context_chunks[0] if context_chunks else ""
        if not top_chunk:
            return ChatResponse(reply="I don't have enough course material yet, please ask your professor to upload content.")
        
        # Clean conversational fallback response
        reply = f"Based on the course materials: {top_chunk[:500]}. For {req.student_level} level students, this covers the essentials. Let me know if you need more details!"
        return ChatResponse(reply=reply)

    # Real RAG + LLM mode.
    retriever = None
    if openai_embeddings_enabled() and course_faiss_exists(course_key):
        embeddings = OpenAIEmbeddings()
        vector_store = FAISS.load_local(course_faiss_path(course_key), embeddings, allow_dangerous_deserialization=True)
        retriever = vector_store.as_retriever(search_kwargs={"k": 4})

    # Build prompt with stored teaching style.
    system_prompt = (
        "You are a friendly teaching assistant. Answer the student's question concisely using the provided course material.\n\n"
        "Rules:\n"
        "1. Answer like a friendly teaching assistant.\n"
        "2. DO NOT include section titles like 'Course Material Response', 'Relevant Material', etc.\n"
        "3. Use a simple explanation based on course material.\n"
        "4. Keep response concise and natural (max 5-6 lines).\n"
        "5. If no relevant material is found in the context, say: 'I don’t have enough course material yet, please ask your professor to upload content.'\n"
        "6. Return ONLY the final answer string. Use plain text and avoid markdown headings.\n"
        f"7. Adaptation: Student level is '{req.student_level}'. Professor's style: '{teaching_style or 'natural'}'.\n\n"
        "Context:\n{context}"
    )

    prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        MessagesPlaceholder(variable_name="chat_history"),
        ("human", "{input}"),
    ])

    llm = ChatGroq(model_name="llama3-8b-8192", temperature=0.7)

    # Generate LangChain memory messages (if any).
    chat_history_messages = []
    for msg in req.history or []:
        if msg.get("sender") == "user":
            chat_history_messages.append(HumanMessage(content=msg.get("text", "")))
        elif msg.get("sender") == "agent":
            chat_history_messages.append(AIMessage(content=msg.get("text", "")))

    if not retriever:
        # Embeddings are disabled; fall back to stuffing the retrieved context directly.
        rag_docs = [{"page_content": c} for c in context_chunks[:4]]
        # Stuff chain expects Document objects; fallback to a simpler template reply using the first chunk.
        return ChatResponse(
            reply=(
                "I found the relevant course context, but I can't run FAISS retrieval in this environment. "
                "I'll still answer using the retrieved snippet.\\n\\n"
                f"Snippet: {context_chunks[0][:600]}..."
            )
        )

    question_answer_chain = create_stuff_documents_chain(llm, prompt)
    rag_chain = create_retrieval_chain(retriever, question_answer_chain)

    response = rag_chain.invoke({
        "input": req.message,
        "chat_history": chat_history_messages,
    })

    return ChatResponse(reply=response.get("answer", ""))

@app.post("/chat/voice", response_model=MultimodalChatResponse)
async def chat_voice(
    course_key: str = Form(...),
    student_level: str = Form("intermediate"),
    message: str = Form(""),
    history: str = Form("[]"),
    audio: UploadFile = File(...),
):
    print(f"[AI Service] /chat/voice | course={course_key} | student_level={student_level} | message_len={len(message)} | history_len={len(history)}")
    
    course_key_norm = normalize_course_key(course_key)
    temp_dir = f"temp_voice_{course_key_norm}_{datetime.utcnow().strftime('%Y%m%d%H%M%S%f')}"
    os.makedirs(temp_dir, exist_ok=True)

    try:
        filename = audio.filename or "voice.webm"
        safe_name = re.sub(r"[^a-zA-Z0-9._-]+", "_", filename)
        file_location = os.path.join(temp_dir, safe_name)
        with open(file_location, "wb") as buffer:
            shutil.copyfileobj(audio.file, buffer)

        transcript = extract_text_from_file(file_location, filename)
        print(f"[AI Service] Voice Transcription: '{transcript[:100]}...'")
        
        final_message = str(message or transcript or "").strip()
        reply = generate_course_chat_reply(
            course_key=course_key_norm,
            message=final_message or "Transcribe and answer this voice message.",
            history=parse_history_payload(history),
            student_level=student_level or "intermediate",
        )

        return MultimodalChatResponse(
            reply=reply.reply,
            transcript=transcript,
            extracted_text=transcript,
            attachments=[{
                "filename": filename,
                "file_type": "audio",
                "text": transcript,
            }],
        )
    except Exception as e:
        print(f"[AI Service] ERROR in /chat/voice: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Voice processing failed: {str(e)}")
    finally:
        try:
            shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception:
            pass

@app.post("/chat/upload", response_model=MultimodalChatResponse)
async def chat_upload(
    course_key: str = Form(...),
    student_level: str = Form("intermediate"),
    message: str = Form(""),
    history: str = Form("[]"),
    files: List[UploadFile] = File(...),
):
    print(f"[AI Service] /chat/upload | course={course_key} | files={len(files)} | message='{message}'")
    
    course_key_norm = normalize_course_key(course_key)
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")

    temp_dir = f"temp_upload_{course_key_norm}_{datetime.utcnow().strftime('%Y%m%d%H%M%S%f')}"
    os.makedirs(temp_dir, exist_ok=True)

    try:
        attachment_records = extract_attachment_records(files, temp_dir)
        attachment_context = build_attachment_context(attachment_records)
        print(f"[AI Service] Extracted attachment context length: {len(attachment_context)}")
        
        prompt = str(message or "Analyze the attached file(s) in course context.").strip()
        reply = generate_course_chat_reply(
            course_key=course_key_norm,
            message=prompt,
            history=parse_history_payload(history),
            student_level=student_level or "intermediate",
            attachment_context=attachment_context,
        )

        extracted_text = "\n\n".join([record.get("text", "") for record in attachment_records if record.get("text")]).strip()
        return MultimodalChatResponse(
            reply=reply.reply,
            transcript="",
            extracted_text=extracted_text,
            attachments=[{
                "filename": record.get("filename", "file"),
                "file_type": record.get("file_type", "document"),
                "text": record.get("text", ""),
            } for record in attachment_records],
        )
    except Exception as e:
        print(f"[AI Service] ERROR in /chat/upload: {str(e)}")
        raise HTTPException(status_code=500, detail=f"File processing failed: {str(e)}")
    finally:
        try:
            shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception:
            pass

@app.post("/personalize", response_model=PersonalizeResponse)
async def generate_personalization(req: PersonalizeRequest):
    if not LLM_AVAILABLE or os.environ.get("GROQ_API_KEY") == "gsk-placeholder":
        # Mock logic if AI isn't connected
        return PersonalizeResponse(
            next_best_topic="Advanced React Hooks",
            topics_to_revise=req.weak_topics if req.weak_topics else ["State Management"],
            practice_questions=[
                "How does useEffect handle dependencies?",
                "What is the difference between useMemo and useCallback?"
            ],
            adaptive_message=f"Agent Notice: Based on your recent quizzes averaging {sum(req.quiz_scores.values())/max(len(req.quiz_scores), 1)}%, let's focus on foundational concepts before moving fast."
        )

    llm = ChatGroq(model_name="llama3-8b-8192", temperature=0.7)
    
    prompt = ChatPromptTemplate.from_messages([
        ("system", (
            "You are an intelligent Course Agent that generates personalized learning paths.\n\n"
            "Given the student's data:\n"
            "Quiz Scores: {quiz_scores}\n"
            "Weak Topics: {weak_topics}\n"
            "Recent Interactions: {recent_interactions}\n\n"
            "Determine exactly:\n"
            "1. The absolute best 'Next Topic' to learn right now.\n"
            "2. A list of 2-3 topics they must revise.\n"
            "3. 2 practice questions tailored to their weak areas.\n"
            "4. A short, encouraging message outlining their adaptive difficulty scale.\n\n"
            "Return the data directly in this precise format (NO MARKDOWN OR OTHER TEXT):\n"
            "Next Best Topic: [Topic]\n"
            "Revise: [Topic 1, Topic 2]\n"
            "Questions: [Q1|Q2]\n"
            "Message: [Message]"
        ))
    ])
    
    try:
        chain = prompt | llm
        response = chain.invoke({
            "quiz_scores": str(req.quiz_scores),
            "weak_topics": ", ".join(req.weak_topics),
            "recent_interactions": ", ".join(req.recent_interactions)
        })
        
        # Simple text parsing since JSON format can be fragile with 8b models without strict structured output
        text = response.content
        lines = text.split("\n")
        
        next_topic = "React Context API" 
        revise = ["Props passing"]
        questions = ["How do you avoid prop drilling?"]
        msg = "Keep practicing, you're doing great!"
        
        for line in lines:
            if line.startswith("Next Best Topic: "): next_topic = line.replace("Next Best Topic: ", "").strip()
            elif line.startswith("Revise: "): revise = [t.strip() for t in line.replace("Revise: ", "").strip("[]").split(",")]
            elif line.startswith("Questions: "): questions = [q.strip() for q in line.replace("Questions: ", "").strip("[]").split("|")]
            elif line.startswith("Message: "): msg = line.replace("Message: ", "").strip()

        return PersonalizeResponse(
            next_best_topic=next_topic,
            topics_to_revise=revise,
            practice_questions=questions,
            adaptive_message=msg
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/evaluate", response_model=EvalResponse)
async def evaluate_submission(req: EvalRequest):
    try:
        return evaluate_submission_core(req.submission_text, req.rubric)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM parsing failed: {str(e)}")

@app.post("/extract/files")
async def extract_files_text(files: List[UploadFile] = File(...)):
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")

    temp_dir = "temp_extract_upload"
    os.makedirs(temp_dir, exist_ok=True)
    parts = []

    try:
        for uploaded in files:
            filename = uploaded.filename or "uploaded-file"
            safe_name = re.sub(r"[^a-zA-Z0-9._-]+", "_", filename)
            file_location = os.path.join(temp_dir, safe_name)

            with open(file_location, "wb") as buffer:
                shutil.copyfileobj(uploaded.file, buffer)

            text = extract_text_from_file(file_location, filename)
            if text:
                parts.append(f"### File: {filename}\n{text[:25000]}")

        return {
            "file_count": len(files),
            "extracted_text": "\n\n".join(parts).strip(),
        }
    finally:
        try:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception:
            pass

@app.post("/assignment/evaluate", response_model=AssignmentEvalResponse)
async def evaluate_assignment_submission(req: AssignmentEvalRequest):
    assignment_text = str(req.assignment_text or "").strip()
    submission_text = str(req.submission_text or "").strip()

    if not assignment_text:
        raise HTTPException(status_code=400, detail="assignment_text is required")
    if not submission_text:
        raise HTTPException(status_code=400, detail="submission_text is required")

    try:
        return evaluate_assignment_core(assignment_text, submission_text, req.rubric)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Assignment evaluation failed: {str(e)}")

@app.post("/evaluate/files", response_model=EvalResponse)
async def evaluate_submission_files(
    rubric: str = Form(...),
    files: List[UploadFile] = File(...),
):
    if not files:
        raise HTTPException(status_code=400, detail="No submission files provided")

    temp_dir = "temp_eval_upload"
    os.makedirs(temp_dir, exist_ok=True)
    parts = []

    try:
        for uploaded in files:
            filename = uploaded.filename or "submission-file"
            safe_name = re.sub(r"[^a-zA-Z0-9._-]+", "_", filename)
            file_location = os.path.join(temp_dir, safe_name)

            with open(file_location, "wb") as buffer:
                shutil.copyfileobj(uploaded.file, buffer)

            text = extract_text_from_file(file_location, filename)
            if not text:
                # Fallback for plain text/code documents.
                try:
                    with open(file_location, "r", encoding="utf-8", errors="ignore") as f:
                        text = f.read()
                except Exception:
                    text = ""

            if text:
                parts.append(f"### File: {filename}\n{text[:20000]}")

        if not parts:
            raise HTTPException(status_code=400, detail="Could not parse readable content from submission files")

        combined = "\n\n".join(parts)
        return evaluate_submission_core(combined, rubric)
    finally:
        try:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception:
            pass

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
